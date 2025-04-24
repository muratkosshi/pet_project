import os

import markdown
from PIL.ImageQt import rgb
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.slide import Slide
from pptx.text.text import Font
from pptx.util import Inches, Pt

from app.engine.mdtree.ImageSearch import ImageSearch
# Локальные модули (пример)
from app.engine.mdtree.MarkdownCategory import MarkdownCategory
from app.engine.mdtree.utils import get_random_file


class MD2Slide:
    """
    Класс, формирующий один слайд в презентации на основе:
      - title (заголовок)
      - content (основное содержимое, часто с <p>...</p> тегами)
      - theme_path (папка с фоновыми изображениями)
      - prompt (для генератора изображений)
      - настройки шрифта (font_name, font_title_size, font_content_size, font_title_color, font_content_color)
    """
    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    theme_title: str = None
    uuid: str = None
    font_name: str = "Times New Roman"
    font_title_size: Pt = Pt(50)
    font_content_size: Pt = Pt(28)
    font_title_color: rgb = RGBColor(0, 0, 0)
    font_content_color: rgb = RGBColor(0, 0, 0)

    def __init__(
        self,
        presentation,
        theme_path,
        title,
        theme_title,
        content,
        prompt,
        *args,
        uuid,
        **kwargs
    ):
        self.presentation = presentation
        # Добавляем слайд на основе макета (layout) №8
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[8])

        # Сохраняем параметры
        self.title = title
        self.theme_title = theme_title
        self.prompt = prompt
        self.content = content
        self.theme = theme_path
        self.uuid = uuid

        # Инициализируем шрифты и создаём слайд
        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()
        self.init_content()

    def imageAdd(self, title, temp_folder_path, theme_title, prompt):
        """
        Ищет или генерирует изображение по теме (title, theme_title, prompt)
        через ImageSearch и добавляет его на слайд. Затем удаляет файл.
        """
        image_search = ImageSearch(
            title,
            theme_title,
            temp_folder_path,
            generate=True,
            prompt=prompt,
            own_domain="https://869d-34-16-183-124.ngrok-free.app"
        )
        image_path = image_search.image_path

        if image_path and os.path.isfile(image_path):
            slide_width = self.presentation.slide_width
            slide_height = self.presentation.slide_height
            print("Adding image:", image_path, "to slide at size:", slide_width, "x", slide_height)

            # Добавляем картинку (примерно в правую часть слайда)
            self.slide.shapes.add_picture(
                image_path,
                left=12200000,    # координаты
                top=0,
                width=slide_width / 2.5,
                height=slide_height
            )

            # Удаляем временный файл изображения
            os.remove(image_path)
        else:
            print(f"Error: {image_path} is not a regular file or doesn't exist.")

    def init_slide(self):
        """
        Добавляет сгенерированное изображение (imageAdd) и ставит фоновое.
        """
        temp_folder_path = f"./temp_images/{self.uuid}"
        if not os.path.exists(temp_folder_path):
            os.makedirs(temp_folder_path)

        # Пытаемся добавить картинку
        try:
            self.imageAdd(self.title, temp_folder_path, self.theme_title, self.prompt)
        except Exception as e:
            print('Не удалось добавить изображение:', e)

        # Меняем фон через placeholder (layout[8] предполагает,
        # что placeholders[1] используется под фоновое изображение)
        placeholder1 = self.slide.placeholders[1]
        path = get_random_file(self.theme)  # получаем случайное изображение из папки
        picture = placeholder1.insert_picture(path)

        # Удаляем placeholder2, чтобы не мешал
        placeholder2 = self.slide.placeholders[2]
        placeholder2.element.getparent().remove(placeholder2.element)

        # Растягиваем картинку на весь слайд
        picture.left = 0
        picture.top = 0
        picture.width = self.presentation.slide_width
        picture.height = self.presentation.slide_height

    def init_font(self, **kwargs):
        """
        Инициализация параметров шрифтов (названия, размеры, цвета)
        на основе переданных аргументов или значений по умолчанию.
        """
        if 'font_name' in kwargs:
            self.font_name = kwargs['font_name']
        if 'font_title_size' in kwargs:
            self.font_title_size = kwargs['font_title_size']
        if 'font_content_size' in kwargs:
            self.font_content_size = kwargs['font_content_size']
        if 'font_title_color' in kwargs:
            self.font_title_color = kwargs['font_title_color']
        if 'font_content_color' in kwargs:
            self.font_content_color = kwargs['font_content_color']

    def get_font(self, font: Font, category: str):
        """
        Применяет настройки шрифта (жирность, имя, размер, цвет)
        в зависимости от категории (TITLE или CONTENT).
        """
        font.bold = True
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
        elif category == MarkdownCategory.CONTENT:
            font.size = self.font_content_size
            font.color.rgb = self.font_content_color

    def init_title(self):
        """
        Создаёт текстовый блок под заголовок (title).
        """
        shapes = self.slide.shapes

        # Задаём ширину для блока с заголовком
        text_box_width = Inches(8)
        # Рассчитываем начальные координаты,
        # чтобы примерно по центру (по горизонтали) отобразить
        start_x = (self.presentation.slide_width - text_box_width) / 5
        start_y = Inches(0.3)

        # Добавляем textbox
        text_box = shapes.add_textbox(start_x, start_y, text_box_width, Inches(0.8))
        tf = text_box.text_frame
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        # Собственно заголовок
        paragraph = tf.paragraphs[0]
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.word_wrap = True

    def init_content(self):
        """
        Создаёт текстовый блок для основного содержимого (content).
        Убираем <p> и </p>, заменяя на переносы строк, затем выводим.
        """
        shapes = self.slide.shapes
        text_box_content = shapes.add_textbox(Inches(2), Inches(3), Inches(10), Inches(7))
        tf = text_box_content.text_frame


        paragraph = tf.paragraphs[0]
        # Убираем теги <p>...</p>, подменяя их на переносы строк
        clean_content = self.content.replace("<p>", "").replace("</p>", "\n")
        paragraph.text = clean_content

        # При желании, можем посмотреть, как markdown-парсер конвертирует это в HTML
        self.processing_md_str(clean_content)

        # Применяем шрифт для контента
        self.get_font(paragraph.font, MarkdownCategory.CONTENT)
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True

    def processing_md_str(self, md_str):
        """
        Пример простого использования markdown в python:
        выводим md_str в консоль, затем печатаем результат конвертации в HTML.
        (Можно дальше при желании доп. парсить и вставлять в pptx как списки и т.д.)
        """
        print("Raw content (for debug):", md_str)
        md = markdown.Markdown()
        html_result = md.convert(md_str)
        print("HTML result:", html_result)
